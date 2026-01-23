# ============================================================
# Presentation Service - Professional PPTX Export Service
# 专业级 PPTX 导出服务 (基于 pptx-skills 设计理念)
# ============================================================

import io
import re
import requests
from typing import Dict, Any, List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from lxml import etree

from app.services.layout_engine import LayoutType


class ProfessionalTypography:
    """专业排版参数 - 参考 pptx-skills 设计标准"""

    # 字体大小 (pt) - 建立清晰的视觉层次
    COVER_TITLE = Pt(48)          # 封面标题 - 最大最醒目
    SECTION_TITLE = Pt(44)        # 章节标题
    SLIDE_TITLE = Pt(36)          # 幻灯片标题
    SUBTITLE = Pt(24)             # 副标题
    BODY = Pt(20)                 # 正文
    BODY_SMALL = Pt(18)           # 小号正文
    CAPTION = Pt(14)              # 说明文字
    QUOTE = Pt(28)                # 引用文字

    # 行间距 - 提升可读性
    LINE_SPACING = 1.5            # 1.5 倍行高
    LINE_SPACING_TIGHT = 1.2      # 紧凑行高

    # 段落间距 (pt)
    PARA_SPACE_BEFORE = Pt(6)
    PARA_SPACE_AFTER = Pt(12)
    BULLET_SPACE_AFTER = Pt(8)    # 列表项间距

    # 字体族
    FONT_TITLE = "Microsoft YaHei"      # 标题字体 (中文)
    FONT_TITLE_EN = "Arial"             # 标题字体 (英文)
    FONT_BODY = "Microsoft YaHei"       # 正文字体
    FONT_BODY_EN = "Arial"              # 正文字体 (英文)


class PptxExportService:
    """
    专业级 PPTX 导出服务
    基于 pptx-skills 设计理念，使用 python-pptx 生成商用级 PowerPoint 文件
    """

    # 幻灯片尺寸 (16:9) - 标准宽屏比例
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # 边距系统 - 专业级间距
    MARGIN = Inches(0.6)                    # 标准边距
    MARGIN_NARROW = Inches(0.4)             # 窄边距
    MARGIN_WIDE = Inches(0.8)               # 宽边距
    CONTENT_WIDTH = Inches(12.133)          # 内容区宽度
    CONTENT_HEIGHT = Inches(6.3)            # 内容区高度

    # 装饰线参数
    DECOR_LINE_WIDTH = Inches(2.5)          # 装饰线宽度
    DECOR_LINE_HEIGHT = Pt(4)               # 装饰线高度
    DECOR_LINE_SHORT = Inches(1.5)          # 短装饰线

    # 专业配色系统 - 18 种主题 (参考 pptx-skills)
    THEME_COLORS = {
        "modern_business": {
            "background": "FFFFFF", "title": "1E3A8A", "subtitle": "64748B",
            "text": "1E293B", "accent": "3B82F6", "accent2": "60A5FA",
            "decorLine": "3B82F6", "bullet": "3B82F6"
        },
        "corporate_blue": {
            "background": "FFFFFF", "title": "1E40AF", "subtitle": "6B7280",
            "text": "1F2937", "accent": "2563EB", "accent2": "3B82F6",
            "decorLine": "2563EB", "bullet": "2563EB"
        },
        "classic_blue": {
            "background": "F4F6F6", "title": "1C2833", "subtitle": "2E4053",
            "text": "2E4053", "accent": "1C2833", "accent2": "AAB7B8",
            "decorLine": "1C2833", "bullet": "1C2833"
        },
        "teal_coral": {
            "background": "FFFFFF", "title": "277884", "subtitle": "5EA8A7",
            "text": "2C3E50", "accent": "FE4447", "accent2": "5EA8A7",
            "decorLine": "FE4447", "bullet": "5EA8A7"
        },
        "elegant_dark": {
            "background": "1A1A1A", "title": "D4AF37", "subtitle": "A0A0A0",
            "text": "F4E4BC", "accent": "D4AF37", "accent2": "C0A030",
            "decorLine": "D4AF37", "bullet": "D4AF37"
        },
        "dark_tech": {
            "background": "0A0A0A", "title": "00FF88", "subtitle": "888888",
            "text": "E0E0E0", "accent": "00D4FF", "accent2": "00FF88",
            "decorLine": "00FF88", "bullet": "00D4FF"
        },
        "gradient_purple": {
            "background": "1A1A2E", "title": "E94560", "subtitle": "A0A0A0",
            "text": "EAEAEA", "accent": "E94560", "accent2": "9B59B6",
            "decorLine": "E94560", "bullet": "E94560"
        },
        "neon_future": {
            "background": "0D0D0D", "title": "FF00FF", "subtitle": "888888",
            "text": "FFFFFF", "accent": "00FFFF", "accent2": "FF00FF",
            "decorLine": "FF00FF", "bullet": "00FFFF"
        },
        "minimal_white": {
            "background": "FFFFFF", "title": "111111", "subtitle": "888888",
            "text": "333333", "accent": "666666", "accent2": "999999",
            "decorLine": "333333", "bullet": "666666"
        },
        "nature_green": {
            "background": "F0FDF4", "title": "166534", "subtitle": "6B7280",
            "text": "1F2937", "accent": "22C55E", "accent2": "16A34A",
            "decorLine": "22C55E", "bullet": "22C55E"
        },
        "soft_pastel": {
            "background": "FDF2F8", "title": "BE185D", "subtitle": "9CA3AF",
            "text": "4A4A4A", "accent": "EC4899", "accent2": "F472B6",
            "decorLine": "EC4899", "bullet": "EC4899"
        },
        "creative_colorful": {
            "background": "FFFFFF", "title": "7C3AED", "subtitle": "6B7280",
            "text": "1F2937", "accent": "F59E0B", "accent2": "7C3AED",
            "decorLine": "7C3AED", "bullet": "F59E0B"
        },
        "warm_sunset": {
            "background": "FFFBEB", "title": "C2410C", "subtitle": "78716C",
            "text": "1F2937", "accent": "F97316", "accent2": "FB923C",
            "decorLine": "F97316", "bullet": "F97316"
        },
        "academic_classic": {
            "background": "FFFEF5", "title": "1E3A5F", "subtitle": "7F8C8D",
            "text": "2C3E50", "accent": "8B4513", "accent2": "1E3A5F",
            "decorLine": "8B4513", "bullet": "1E3A5F"
        },
        "anime_dark": {
            "background": "1A1A2E", "title": "FF6B9D", "subtitle": "A78BFA",
            "text": "E0E0E0", "accent": "C084FC", "accent2": "FF6B9D",
            "decorLine": "FF6B9D", "bullet": "C084FC"
        },
        "anime_cute": {
            "background": "FFF0F5", "title": "FF69B4", "subtitle": "DDA0DD",
            "text": "4A4A4A", "accent": "FFB6C1", "accent2": "FF69B4",
            "decorLine": "FF69B4", "bullet": "FFB6C1"
        },
        "cyberpunk": {
            "background": "0D0221", "title": "FF00FF", "subtitle": "FF6EC7",
            "text": "E0E0E0", "accent": "00FFFF", "accent2": "FF00FF",
            "decorLine": "FF00FF", "bullet": "00FFFF"
        },
        "eva_nerv": {
            "background": "1C1C1C", "title": "5B2C6F", "subtitle": "E74C3C",
            "text": "E0E0E0", "accent": "1ABC9C", "accent2": "5B2C6F",
            "decorLine": "5B2C6F", "bullet": "1ABC9C"
        },
        "retro_pixel": {
            "background": "2D1B69", "title": "FF6B6B", "subtitle": "FFE66D",
            "text": "F8F8F8", "accent": "4ECDC4", "accent2": "FF6B6B",
            "decorLine": "FF6B6B", "bullet": "4ECDC4"
        },
    }

    # 排版参数引用
    typo = ProfessionalTypography

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """十六进制转 RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    def _get_theme_colors(self, theme: str) -> Dict[str, str]:
        """获取主题配色"""
        return self.THEME_COLORS.get(theme, self.THEME_COLORS["modern_business"])

    def _set_paragraph_spacing(self, paragraph, space_before: Pt = None, space_after: Pt = None, line_spacing: float = None):
        """设置段落间距"""
        if space_before:
            paragraph.space_before = space_before
        if space_after:
            paragraph.space_after = space_after
        if line_spacing:
            paragraph.line_spacing = line_spacing

    def _add_decoration_line(self, slide, x: Inches, y: Inches, width: Inches, color: str, height: Pt = None):
        """添加装饰线"""
        h = height or self.DECOR_LINE_HEIGHT
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, h)
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(color)
        line.line.fill.background()
        return line

    def _create_text_frame(self, shape, word_wrap: bool = True, margin_left: Inches = None,
                           margin_right: Inches = None, margin_top: Inches = None, margin_bottom: Inches = None):
        """配置文本框属性"""
        tf = shape.text_frame
        tf.word_wrap = word_wrap
        if margin_left is not None:
            tf.margin_left = margin_left
        if margin_right is not None:
            tf.margin_right = margin_right
        if margin_top is not None:
            tf.margin_top = margin_top
        if margin_bottom is not None:
            tf.margin_bottom = margin_bottom
        return tf

    def _parse_content(self, content: str) -> List[Tuple[str, int, bool]]:
        """解析内容为 (文本, 级别, 是否列表项)"""
        if not content:
            return []
        content = content.replace("\\n", "\n")
        lines = []
        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                continue
            level = min(2, (len(line) - len(line.lstrip())) // 2)
            is_bullet = stripped.startswith(("- ", "* ", "+ ")) or bool(re.match(r'^\d+\.\s', stripped))
            text = re.sub(r'^[-*+]\s|^\d+\.\s|^#{1,3}\s', '', stripped)
            lines.append((text, level, is_bullet))
        return lines

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str, colors: Dict[str, str]) -> None:
        """
        封面页 - 专业级设计
        - 垂直居中的大标题
        - 优雅的装饰线
        - 副标题
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self._hex_to_rgb(colors["background"])

        # 标题 - 垂直居中偏上
        title_top = Inches(2.8)
        title_box = slide.shapes.add_textbox(self.MARGIN, title_top, self.CONTENT_WIDTH, Inches(1.2))
        tf = self._create_text_frame(title_box)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self.typo.COVER_TITLE
        p.font.bold = True
        p.font.name = self.typo.FONT_TITLE
        p.font.color.rgb = self._hex_to_rgb(colors["title"])
        p.alignment = PP_ALIGN.CENTER

        # 装饰线 - 居中
        line_x = (self.SLIDE_WIDTH - self.DECOR_LINE_WIDTH) / 2
        self._add_decoration_line(slide, line_x, Inches(4.1), self.DECOR_LINE_WIDTH, colors["decorLine"])

        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(self.MARGIN, Inches(4.4), self.CONTENT_WIDTH, Inches(0.8))
            tf = self._create_text_frame(sub_box)
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = self.typo.SUBTITLE
            p.font.name = self.typo.FONT_BODY
            p.font.color.rgb = self._hex_to_rgb(colors["subtitle"])
            p.alignment = PP_ALIGN.CENTER

    def _add_section_slide(self, prs: Presentation, title: str, colors: Dict[str, str]) -> None:
        """
        章节页 - 专业级设计
        - 大号章节标题
        - 居中装饰线
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self._hex_to_rgb(colors["background"])

        # 章节标题
        title_box = slide.shapes.add_textbox(self.MARGIN, Inches(3.0), self.CONTENT_WIDTH, Inches(1.2))
        tf = self._create_text_frame(title_box)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self.typo.SECTION_TITLE
        p.font.bold = True
        p.font.name = self.typo.FONT_TITLE
        p.font.color.rgb = self._hex_to_rgb(colors["title"])
        p.alignment = PP_ALIGN.CENTER

        # 装饰线
        line_x = (self.SLIDE_WIDTH - self.DECOR_LINE_WIDTH) / 2
        self._add_decoration_line(slide, line_x, Inches(4.3), self.DECOR_LINE_WIDTH, colors["decorLine"])

    def _add_content_slide(self, prs: Presentation, title: str, content: str, colors: Dict[str, str], notes: str = "") -> None:
        """
        内容页 - 专业级设计
        - 紧凑的标题区
        - 短装饰线
        - 内容感知的动态布局
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self._hex_to_rgb(colors["background"])

        # 标题 - 紧凑布局，减少顶部空白
        title_box = slide.shapes.add_textbox(self.MARGIN, Inches(0.4), self.CONTENT_WIDTH, Inches(0.7))
        tf = self._create_text_frame(title_box)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self.typo.SLIDE_TITLE
        p.font.bold = True
        p.font.name = self.typo.FONT_TITLE
        p.font.color.rgb = self._hex_to_rgb(colors["title"])

        # 短装饰线 - 紧跟标题
        self._add_decoration_line(slide, self.MARGIN, Inches(1.1), self.DECOR_LINE_SHORT, colors["decorLine"], Pt(3))

        # 解析内容
        lines = self._parse_content(content)
        num_lines = len(lines)

        # 幻灯片总高度 7.5 英寸，底部留 0.4 英寸边距
        slide_usable_bottom = Inches(7.1)
        content_area_bottom = Inches(1.1)  # 装饰线位置

        # 根据内容量动态计算起始位置
        # 幻灯片中心大约在 3.75 英寸
        slide_center = Inches(3.75)

        if num_lines == 0:
            # 无内容
            content_top = Inches(1.35)
        elif num_lines == 1:
            # 单行内容 - 明显居中（放在中心位置）
            content_top = slide_center - Inches(0.3)  # 稍微偏上一点
        elif num_lines == 2:
            # 两行内容 - 居中
            content_top = slide_center - Inches(0.6)
        elif num_lines == 3:
            # 三行内容 - 轻微居中
            content_top = slide_center - Inches(0.9)
        elif num_lines <= 5:
            # 4-5行 - 轻微偏移
            content_top = Inches(2.0)
        else:
            # 内容多 - 从标准位置开始
            content_top = Inches(1.35)

        content_box = slide.shapes.add_textbox(self.MARGIN, content_top, self.CONTENT_WIDTH, Inches(5.8))
        tf = self._create_text_frame(content_box, margin_left=Inches(0.1), margin_top=Inches(0.05))

        if lines:
            # 第一行
            text, level, is_bullet = lines[0]
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = self.typo.BODY
            p.font.name = self.typo.FONT_BODY
            p.font.color.rgb = self._hex_to_rgb(colors["text"])
            p.level = level
            if is_bullet:
                p.bullet = True
            self._set_paragraph_spacing(p, space_after=self.typo.PARA_SPACE_AFTER)

            # 后续行
            for text, level, is_bullet in lines[1:]:
                p = tf.add_paragraph()
                p.text = text
                p.font.size = Pt(20 - level * 2)
                p.font.name = self.typo.FONT_BODY
                p.font.color.rgb = self._hex_to_rgb(colors["text"])
                p.level = level
                if is_bullet:
                    p.bullet = True
                self._set_paragraph_spacing(p, space_before=Pt(6), space_after=self.typo.PARA_SPACE_AFTER)

        # 演讲者备注
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_two_column_slide(self, prs: Presentation, title: str, content: str, colors: Dict[str, str], notes: str = "") -> None:
        """
        双栏页 - 专业级设计
        - 紧凑标题区
        - 均衡的双栏布局
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self._hex_to_rgb(colors["background"])

        # 标题 - 紧凑布局
        title_box = slide.shapes.add_textbox(self.MARGIN, Inches(0.4), self.CONTENT_WIDTH, Inches(0.7))
        tf = self._create_text_frame(title_box)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self.typo.SLIDE_TITLE
        p.font.bold = True
        p.font.name = self.typo.FONT_TITLE
        p.font.color.rgb = self._hex_to_rgb(colors["title"])

        # 短装饰线
        self._add_decoration_line(slide, self.MARGIN, Inches(1.1), self.DECOR_LINE_SHORT, colors["decorLine"], Pt(3))

        # 分割内容
        lines = self._parse_content(content)
        mid = len(lines) // 2
        left_lines = lines[:mid] if mid else lines
        right_lines = lines[mid:] if mid else []

        # 列宽计算 - 更紧凑的间隙
        col_width = Inches(5.8)
        col_gap = Inches(0.5)
        left_x = self.MARGIN
        right_x = self.MARGIN + col_width + col_gap
        content_top = Inches(1.35)
        content_height = Inches(5.8)

        # 左栏
        left_box = slide.shapes.add_textbox(left_x, content_top, col_width, content_height)
        tf = self._create_text_frame(left_box, margin_left=Inches(0.05))
        if left_lines:
            p = tf.paragraphs[0]
            p.text = left_lines[0][0]
            p.font.size = self.typo.BODY_SMALL
            p.font.name = self.typo.FONT_BODY
            p.font.color.rgb = self._hex_to_rgb(colors["text"])
            if left_lines[0][2]:
                p.bullet = True
            self._set_paragraph_spacing(p, space_after=self.typo.PARA_SPACE_AFTER)

            for text, _, is_bullet in left_lines[1:]:
                p = tf.add_paragraph()
                p.text = text
                p.font.size = self.typo.BODY_SMALL
                p.font.name = self.typo.FONT_BODY
                p.font.color.rgb = self._hex_to_rgb(colors["text"])
                if is_bullet:
                    p.bullet = True
                self._set_paragraph_spacing(p, space_before=Pt(6), space_after=self.typo.PARA_SPACE_AFTER)

        # 右栏
        right_box = slide.shapes.add_textbox(right_x, content_top, col_width, content_height)
        tf = self._create_text_frame(right_box, margin_left=Inches(0.05))
        if right_lines:
            p = tf.paragraphs[0]
            p.text = right_lines[0][0]
            p.font.size = self.typo.BODY_SMALL
            p.font.name = self.typo.FONT_BODY
            p.font.color.rgb = self._hex_to_rgb(colors["text"])
            if right_lines[0][2]:
                p.bullet = True
            self._set_paragraph_spacing(p, space_after=self.typo.PARA_SPACE_AFTER)

            for text, _, is_bullet in right_lines[1:]:
                p = tf.add_paragraph()
                p.text = text
                p.font.size = self.typo.BODY_SMALL
                p.font.name = self.typo.FONT_BODY
                p.font.color.rgb = self._hex_to_rgb(colors["text"])
                if is_bullet:
                    p.bullet = True
                self._set_paragraph_spacing(p, space_before=Pt(6), space_after=self.typo.PARA_SPACE_AFTER)

        # 演讲者备注
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_quote_slide(self, prs: Presentation, title: str, content: str, colors: Dict[str, str], notes: str = "") -> None:
        """
        引用页 - 专业级设计
        - 大引号装饰
        - 斜体引用文字
        - 来源署名
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self._hex_to_rgb(colors["background"])

        # 大引号装饰
        quote_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(1), Inches(1.2))
        tf = self._create_text_frame(quote_box)
        p = tf.paragraphs[0]
        p.text = "\u201C"
        p.font.size = Pt(96)
        p.font.color.rgb = self._hex_to_rgb(colors["accent"])

        # 引用内容
        content_box = slide.shapes.add_textbox(Inches(2.2), Inches(2.5), Inches(9), Inches(2.5))
        tf = self._create_text_frame(content_box)
        p = tf.paragraphs[0]
        p.text = content.replace("\\n", "\n").strip()
        p.font.size = self.typo.QUOTE
        p.font.italic = True
        p.font.name = self.typo.FONT_BODY
        p.font.color.rgb = self._hex_to_rgb(colors["text"])
        p.alignment = PP_ALIGN.CENTER
        self._set_paragraph_spacing(p, line_spacing=1.4)

        # 来源署名
        if title:
            source_box = slide.shapes.add_textbox(Inches(2.2), Inches(5.3), Inches(9), Inches(0.6))
            tf = self._create_text_frame(source_box)
            p = tf.paragraphs[0]
            p.text = f"— {title}"
            p.font.size = self.typo.BODY_SMALL
            p.font.name = self.typo.FONT_BODY
            p.font.color.rgb = self._hex_to_rgb(colors["subtitle"])
            p.alignment = PP_ALIGN.CENTER

        # 演讲者备注
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_thank_you_slide(self, prs: Presentation, title: str, content: str, colors: Dict[str, str]) -> None:
        """
        感谢页 - 专业级设计
        - 装饰线
        - 大标题
        - 联系信息
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self._hex_to_rgb(colors["background"])

        # 上装饰线
        line_x = (self.SLIDE_WIDTH - self.DECOR_LINE_WIDTH) / 2
        self._add_decoration_line(slide, line_x, Inches(2.6), self.DECOR_LINE_WIDTH, colors["decorLine"])

        # 感谢标题
        title_box = slide.shapes.add_textbox(self.MARGIN, Inches(2.9), self.CONTENT_WIDTH, Inches(1.5))
        tf = self._create_text_frame(title_box)
        p = tf.paragraphs[0]
        p.text = title or "Thank You"
        p.font.size = self.typo.COVER_TITLE
        p.font.bold = True
        p.font.name = self.typo.FONT_TITLE
        p.font.color.rgb = self._hex_to_rgb(colors["title"])
        p.alignment = PP_ALIGN.CENTER

        # 下装饰线
        self._add_decoration_line(slide, line_x, Inches(4.4), self.DECOR_LINE_WIDTH, colors["decorLine"])

        # 联系信息/副标题
        if content:
            sub_box = slide.shapes.add_textbox(self.MARGIN, Inches(4.8), self.CONTENT_WIDTH, Inches(1))
            tf = self._create_text_frame(sub_box)
            p = tf.paragraphs[0]
            p.text = content.replace("\\n", " ").strip()
            p.font.size = self.typo.BODY
            p.font.name = self.typo.FONT_BODY
            p.font.color.rgb = self._hex_to_rgb(colors["subtitle"])
            p.alignment = PP_ALIGN.CENTER

    def _create_slide(self, prs: Presentation, slide_data: Dict[str, Any], colors: Dict[str, str]) -> None:
        """根据布局创建幻灯片"""
        layout = slide_data.get("layout", "bullet_points")
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        notes = slide_data.get("notes", "")

        if layout == LayoutType.TITLE_COVER.value:
            self._add_title_slide(prs, title, content, colors)
        elif layout == LayoutType.TITLE_SECTION.value:
            self._add_section_slide(prs, title, colors)
        elif layout == LayoutType.QUOTE_CENTER.value:
            self._add_quote_slide(prs, title, content, colors, notes)
        elif layout == LayoutType.THANK_YOU.value:
            self._add_thank_you_slide(prs, title, content, colors)
        elif layout in (LayoutType.TWO_COLUMN.value, LayoutType.THREE_COLUMN.value, LayoutType.COMPARISON.value):
            self._add_two_column_slide(prs, title, content, colors, notes)
        else:
            self._add_content_slide(prs, title, content, colors, notes)

    async def export_to_pptx(self, presentation_data: Dict[str, Any], theme: str = "modern_business") -> bytes:
        """导出为 PPTX"""
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT
        colors = self._get_theme_colors(theme)

        for slide_data in presentation_data.get("slides", []):
            self._create_slide(prs, slide_data, colors)

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()

    def generate_filename(self, title: str) -> str:
        """生成文件名"""
        from datetime import datetime
        safe_title = "".join(c for c in title.lower().replace(" ", "_") if c.isalnum() or c in "_-")[:50]
        return f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"


# 全局实例
pptx_export_service = PptxExportService()
