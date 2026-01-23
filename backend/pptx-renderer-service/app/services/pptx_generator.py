# ============================================================
# PPTX Generator Service - PPTX 生成服务
# ============================================================
"""
根据元素位置信息生成 PPTX 文件
"""

import io
import logging
import re
from typing import List, Dict, Any, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import httpx

from app.config import settings
from app.services.position_extractor import ElementPosition

logger = logging.getLogger(__name__)


class PptxGeneratorService:
    """
    PPTX 生成服务

    根据从 HTML 提取的元素位置信息生成 PPTX 文件。
    支持文本、图片、形状等元素类型。
    """

    def __init__(
        self,
        slide_width_inches: float = None,
        slide_height_inches: float = None,
    ):
        """
        初始化 PPTX 生成器

        Args:
            slide_width_inches: 幻灯片宽度 (英寸)
            slide_height_inches: 幻灯片高度 (英寸)
        """
        self.slide_width_inches = slide_width_inches or settings.PPTX_SLIDE_WIDTH_INCHES
        self.slide_height_inches = slide_height_inches or settings.PPTX_SLIDE_HEIGHT_INCHES

        # 像素到英寸的转换比例 (基于 1920x1080 分辨率)
        self.px_to_inch_x = self.slide_width_inches / 1920
        self.px_to_inch_y = self.slide_height_inches / 1080

    def generate(
        self,
        elements: List[ElementPosition],
        screenshot: bytes = None,
        render_width: int = 1920,
        render_height: int = 1080,
        use_screenshot_background: bool = False,
    ) -> bytes:
        """
        根据元素位置生成 PPTX

        Args:
            elements: 元素位置列表
            screenshot: 可选的截图作为背景
            render_width: 渲染宽度 (用于坐标转换)
            render_height: 渲染高度 (用于坐标转换)
            use_screenshot_background: 是否使用截图作为背景

        Returns:
            bytes: PPTX 文件内容
        """
        # 更新转换比例
        self.px_to_inch_x = self.slide_width_inches / render_width
        self.px_to_inch_y = self.slide_height_inches / render_height

        # 创建演示文稿
        prs = Presentation()

        # 设置幻灯片尺寸
        prs.slide_width = Inches(self.slide_width_inches)
        prs.slide_height = Inches(self.slide_height_inches)

        # 添加空白幻灯片
        blank_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(blank_layout)

        # 如果使用截图作为背景
        if use_screenshot_background and screenshot:
            self._add_background_image(slide, screenshot)
        else:
            # 根据元素生成内容
            for element in elements:
                try:
                    self._add_element(slide, element)
                except Exception as e:
                    logger.warning(f"添加元素失败: {element.element_id}, 错误: {e}")

        # 保存到字节流
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        logger.info(f"PPTX 生成完成: {len(elements)} 个元素")
        return output.read()

    def generate_multi_slide(
        self,
        slides_data: List[Dict[str, Any]],
    ) -> bytes:
        """
        生成多页幻灯片

        Args:
            slides_data: 幻灯片数据列表，每项包含 elements 和可选的 screenshot

        Returns:
            bytes: PPTX 文件内容
        """
        prs = Presentation()
        prs.slide_width = Inches(self.slide_width_inches)
        prs.slide_height = Inches(self.slide_height_inches)

        blank_layout = prs.slide_layouts[6]

        for i, slide_data in enumerate(slides_data):
            slide = prs.slides.add_slide(blank_layout)
            elements = slide_data.get("elements", [])
            screenshot = slide_data.get("screenshot")
            use_bg = slide_data.get("use_screenshot_background", False)

            if use_bg and screenshot:
                self._add_background_image(slide, screenshot)
            else:
                for element in elements:
                    try:
                        self._add_element(slide, element)
                    except Exception as e:
                        logger.warning(f"幻灯片 {i + 1} 添加元素失败: {e}")

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        logger.info(f"多页 PPTX 生成完成: {len(slides_data)} 页")
        return output.read()

    def _add_element(self, slide, element: ElementPosition):
        """
        添加单个元素到幻灯片

        Args:
            slide: PPTX 幻灯片对象
            element: 元素位置信息
        """
        element_type = element.element_type.lower()

        if element_type in ["text", "title", "subtitle", "list"]:
            self._add_text_element(slide, element)
        elif element_type == "image":
            self._add_image_element(slide, element)
        elif element_type == "shape":
            self._add_shape_element(slide, element)
        else:
            # 默认作为文本处理
            if element.content:
                self._add_text_element(slide, element)

    def _add_text_element(self, slide, element: ElementPosition):
        """添加文本元素"""
        if not element.content:
            return

        # 转换坐标
        left, top, width, height = self._convert_position(element)

        # 创建文本框
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        # 设置文本
        p = tf.paragraphs[0]
        p.text = element.content

        # 设置字体大小
        if element.font_size:
            # 将像素转换为磅 (大约 0.75 倍)
            font_size_pt = element.font_size * 0.75
            p.font.size = Pt(font_size_pt)

        # 设置字体颜色
        if element.color:
            rgb = self._parse_color(element.color)
            if rgb:
                p.font.color.rgb = rgb

        # 设置字体粗细
        if element.font_weight:
            weight = element.font_weight
            if weight == "bold" or (weight.isdigit() and int(weight) >= 700):
                p.font.bold = True

        # 设置对齐方式
        if element.text_align:
            align_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }
            p.alignment = align_map.get(element.text_align, PP_ALIGN.LEFT)

    def _add_image_element(self, slide, element: ElementPosition):
        """添加图片元素"""
        if not element.src:
            return

        # 转换坐标
        left, top, width, height = self._convert_position(element)

        try:
            # 获取图片数据
            image_data = self._fetch_image(element.src)
            if image_data:
                image_stream = io.BytesIO(image_data)
                slide.shapes.add_picture(image_stream, left, top, width, height)
        except Exception as e:
            logger.warning(f"添加图片失败: {element.src}, 错误: {e}")
            # 添加占位符矩形
            self._add_placeholder_shape(slide, left, top, width, height)

    def _add_shape_element(self, slide, element: ElementPosition):
        """添加形状元素"""
        left, top, width, height = self._convert_position(element)

        # 获取形状类型
        shape_type = MSO_SHAPE.RECTANGLE
        if element.attributes:
            shape_name = element.attributes.get("shape", "rectangle").lower()
            shape_map = {
                "rectangle": MSO_SHAPE.RECTANGLE,
                "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
                "oval": MSO_SHAPE.OVAL,
                "circle": MSO_SHAPE.OVAL,
                "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            }
            shape_type = shape_map.get(shape_name, MSO_SHAPE.RECTANGLE)

        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        # 设置填充颜色
        if element.background_color:
            rgb = self._parse_color(element.background_color)
            if rgb:
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb

    def _add_background_image(self, slide, screenshot: bytes):
        """添加背景图片"""
        image_stream = io.BytesIO(screenshot)
        slide.shapes.add_picture(
            image_stream,
            Inches(0),
            Inches(0),
            Inches(self.slide_width_inches),
            Inches(self.slide_height_inches),
        )

    def _add_placeholder_shape(self, slide, left, top, width, height):
        """添加占位符形状 (用于图片加载失败时)"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(200, 200, 200)

    def _convert_position(
        self,
        element: ElementPosition,
    ) -> Tuple[Emu, Emu, Emu, Emu]:
        """
        将像素坐标转换为 PPTX 坐标 (EMU)

        Args:
            element: 元素位置信息

        Returns:
            Tuple[Emu, Emu, Emu, Emu]: (left, top, width, height)
        """
        left = Inches(element.x * self.px_to_inch_x)
        top = Inches(element.y * self.px_to_inch_y)
        width = Inches(element.width * self.px_to_inch_x)
        height = Inches(element.height * self.px_to_inch_y)

        return left, top, width, height

    def _parse_color(self, color_str: str) -> Optional[RGBColor]:
        """
        解析 CSS 颜色字符串为 RGBColor

        Args:
            color_str: CSS 颜色字符串 (如 "rgb(255, 0, 0)" 或 "#ff0000")

        Returns:
            Optional[RGBColor]: RGB 颜色对象
        """
        if not color_str:
            return None

        try:
            # 处理 rgb/rgba 格式
            rgb_match = re.match(
                r"rgba?\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)",
                color_str
            )
            if rgb_match:
                r, g, b = map(int, rgb_match.groups())
                return RGBColor(r, g, b)

            # 处理十六进制格式
            hex_match = re.match(r"#([0-9a-fA-F]{6})", color_str)
            if hex_match:
                hex_color = hex_match.group(1)
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)

            # 处理简写十六进制格式
            hex_short_match = re.match(r"#([0-9a-fA-F]{3})$", color_str)
            if hex_short_match:
                hex_color = hex_short_match.group(1)
                r = int(hex_color[0] * 2, 16)
                g = int(hex_color[1] * 2, 16)
                b = int(hex_color[2] * 2, 16)
                return RGBColor(r, g, b)

        except Exception as e:
            logger.warning(f"解析颜色失败: {color_str}, 错误: {e}")

        return None

    def _fetch_image(self, src: str) -> Optional[bytes]:
        """
        获取图片数据

        Args:
            src: 图片 URL 或 base64 数据

        Returns:
            Optional[bytes]: 图片数据
        """
        if not src:
            return None

        # 处理 base64 数据
        if src.startswith("data:"):
            try:
                # 提取 base64 部分
                base64_data = src.split(",", 1)[1]
                import base64
                return base64.b64decode(base64_data)
            except Exception as e:
                logger.warning(f"解析 base64 图片失败: {e}")
                return None

        # 处理 URL
        if src.startswith(("http://", "https://")):
            try:
                with httpx.Client(timeout=10.0) as client:
                    response = client.get(src)
                    response.raise_for_status()
                    return response.content
            except Exception as e:
                logger.warning(f"下载图片失败: {src}, 错误: {e}")
                return None

        return None
